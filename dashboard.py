import streamlit as st
import pandas as pd
import docx2txt
import json
import PyPDF2
import pptx
from textblob import TextBlob

# ---------------------------
# Helper functions
# ---------------------------
def read_file(file):
    text = ""

    if file.name.endswith(".txt"):
        try:
            text = file.read().decode("utf-8")
        except UnicodeDecodeError:
            file.seek(0)
            text = file.read().decode("latin1", errors="ignore")

    elif file.name.endswith(".pdf"):
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            page_text = page.extract_text() or ""
            text += page_text

    elif file.name.endswith(".docx"):
        text = docx2txt.process(file) or ""

    elif file.name.endswith(".json"):
        try:
            data = json.load(file)
            text = json.dumps(data, indent=2, ensure_ascii=False)
        except UnicodeDecodeError:
            file.seek(0)
            data = json.load(file)
            text = json.dumps(data, indent=2, ensure_ascii=False)

    elif file.name.endswith(".pptx"):
        prs = pptx.Presentation(file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += (shape.text or "") + "\n"

    elif file.name.endswith(".csv"):
        try:
            df = pd.read_csv(file, encoding="utf-8")
        except UnicodeDecodeError:
            file.seek(0)
            df = pd.read_csv(file, encoding="latin1")
        text = df.to_string(index=False)

    elif file.name.endswith(".xls") or file.name.endswith(".xlsx"):
        df = pd.read_excel(file)
        text = df.to_string(index=False)

    else:
        text = "Unsupported file format."
    return text


def analyze_text(text):
    words = text.split()
    word_count = len(words)
    unique_words = len(set(words))
    char_count = len(text)
    reading_time = round(word_count / 200, 2)

    blob = TextBlob(text)
    sentiment = blob.sentiment.polarity
    positivity = round((sentiment + 1) * 50, 2)

    keywords = pd.Series([w.lower() for w in words]).value_counts().head(10).to_dict() if words else {}

    upper_count = sum(1 for ch in text if ch.isalpha() and ch.isupper())
    lower_count = sum(1 for ch in text if ch.isalpha() and ch.islower())

    return {
        "word_count": word_count,
        "unique_words": unique_words,
        "char_count": char_count,
        "reading_time": reading_time,
        "sentiment": sentiment,
        "positivity": positivity,
        "keywords": keywords,
        "upper_count": upper_count,
        "lower_count": lower_count,
    }

# ---------------------------
# Streamlit UI
# ---------------------------
st.set_page_config(page_title="Smart File Analyzer", layout="wide")

# Pink Header
st.markdown(
    """
    <div style="background-color:#ff69b4;padding:15px;border-radius:10px;margin-bottom:20px;">
        <h2 style="color:white;text-align:center;"> Smart File Analyzer Dashboard</h2>
    </div>
    """,
    unsafe_allow_html=True
)

# ---------------------------
# Top Navigation Buttons
# ---------------------------
colA, colB, colC = st.columns([1,1,1])
with colA:
    dashboard_btn = st.button("🏠 Dashboard")
with colB:
    upload_btn = st.button("📂 Upload & Analyze")
with colC:
    about_btn = st.button("ℹ️ About")

# Session state for navigation
if "page" not in st.session_state:
    st.session_state.page = "Dashboard"

if dashboard_btn:
    st.session_state.page = "Dashboard"
elif upload_btn:
    st.session_state.page = "Upload"
elif about_btn:
    st.session_state.page = "About"

# ---------------------------
# Pages
# ---------------------------
if st.session_state.page == "Dashboard":
    st.title("📌 Overview")
    col1, col2, col3 = st.columns(3)
    col1.metric("Files Processed", "12")
    col2.metric("Avg. Sentiment", "0.25 😊")
    col3.metric("Most Common Word", "data")
    st.divider()
    st.subheader("⚡ System Activity")
    st.info("Upload a file in 'Upload & Analyze' to get live insights.")

elif st.session_state.page == "Upload":
    st.title("📂 Upload & Analyze")
    uploaded_file = st.file_uploader(
        "Upload any text-based file",
        type=["txt", "pdf", "docx", "json", "pptx", "csv", "xls", "xlsx"]
    )

    if uploaded_file is not None:
        text = read_file(uploaded_file)

        st.subheader("📑 Extracted Content")
        st.text_area("Preview", text[:1000] + "..." if len(text) > 1000 else text, height=200)

        insights = analyze_text(text)

        st.subheader("📊 Insights Overview")
        col1, col2, col3 = st.columns(3)
        col1.metric("Total Words", insights["word_count"])
        col2.metric("Sentiment Score", round(insights["sentiment"], 3))
        col3.metric("Positivity %", f"{insights['positivity']}%")

        st.divider()
        st.subheader("📦 Extra Analysis")
        col4, col5, col6 = st.columns(3)
        col4.success(f"Unique Words: {insights['unique_words']}")
        col5.warning(f"Characters: {insights['char_count']}")
        col6.info(f"Reading Time: {insights['reading_time']} min")

        col7, col8 = st.columns(2)
        with col7:
            st.markdown(
                f"""
                <div style="background-color:#ffe4ec;padding:20px;border-radius:15px;text-align:center;">
                    <h3>Uppercase Letters (A–Z)</h3>
                    <h1>{insights['upper_count']}</h1>
                </div>
                """,
                unsafe_allow_html=True
            )
        with col8:
            st.markdown(
                f"""
                <div style="background-color:#ffd1dc;padding:20px;border-radius:15px;text-align:center;">
                    <h3>Lowercase Letters (a–z)</h3>
                    <h1>{insights['lower_count']}</h1>
                </div>
                """,
                unsafe_allow_html=True
            )

        st.divider()
        st.subheader("😊 Sentiment Progress")
        st.progress(int(insights["positivity"]))

        st.subheader("🔡 Case Distribution")
        case_df = pd.DataFrame(
            {"Count": [insights["upper_count"], insights["lower_count"]]},
            index=["Uppercase", "Lowercase"]
        )
        st.bar_chart(case_df)

        st.subheader("🔑 Top Keywords")
        st.bar_chart(pd.DataFrame.from_dict(insights["keywords"], orient="index", columns=["Count"]))

        st.divider()
        st.subheader("⚙️ Quick Actions")
        col9, col10, col11 = st.columns(3)
        col9.download_button("⬇️ Download Report", text, file_name="report.txt")
        col10.download_button(
            "📤 Export Keywords CSV",
            pd.DataFrame(list(insights["keywords"].items()), columns=["Word", "Count"]).to_csv(index=False),
            file_name="keywords.csv"
        )
        col11.download_button("📋 Copy Raw Text", text, file_name="raw_text.txt")

elif st.session_state.page == "About":
    st.title("ℹ️ About This App")
    st.markdown("""
      **Smart File Analyzer**  
    - Supports: All Text Data Formats 
    - Auto-handles encoding issues  
    - Multiple insights: word counts, sentiment, uppercase vs lowercase  
    - Interactive charts & progress bars  
    - Quick export options  
    """)
